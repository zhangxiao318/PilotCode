#!/usr/bin/env python3
"""Generate E2E test results PPT for PilotCode: Qwen3.5-9B vs 30B-A3."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0 if not bullet.startswith("  ") else 1
        p.font.size = Pt(20)
    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_shape.text_frame.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x5F, 0x7A)

    for bullet in left_bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0

    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x8B, 0x45, 0x13)

    for bullet in right_bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0

    return slide


# Slide 1: Title
add_title_slide(
    prs,
    "PilotCode E2E 编程能力验证",
    "Qwen3.5-9B vs Qwen3.5-30B-A3 对比测试\n基于 redpanda-dev 大型分布式项目\n2026-04-28",
)

# Slide 2: Target System & Tasks Overview
add_content_slide(
    prs,
    "目标系统与测试任务概览",
    [
        "目标系统：Redpanda（分布式流处理平台，Kafka 兼容）",
        "项目规模：6,151 个源文件，约 127 万行代码（C++ / Go / Python）",
        "测试环境：本地 vLLM + PilotCode WebSocket 模式",
        "",
        "测试任务列表（共 7 项）：",
        "  1. Python 代码分析 — dev_cluster.py 功能架构",
        "  2. Python 简单修改 — 添加 --version 参数",
        "  3. Python 复杂修改 — stream_until_eof 添加时间戳",
        "  4. C++ 代码分析 — cache_probe.cc metrics 体系",
        "  5. C++ 简单修改 — io_result.cc 提取 to_string()",
        "  6. C++ 复杂修改 — cache_probe.cc 添加 cache_hit_ratio",
        "  7. 多文件功能增强 — 配置验证器 + 摘要打印",
    ],
)

# Slide 3: Model Comparison Overview
add_two_column_slide(
    prs,
    "模型对比：Qwen3.5-9B vs Qwen3.5-30B-A3",
    "Qwen3.5-9B（弱模型）",
    [
        "显存占用：~20 GB（单卡可跑）",
        "上下文窗口：32K",
        "代码生成：易幻觉、Loop Guard 触发率高",
        "工具调用：准确率 ~60%",
        "FileEdit：多次尝试、频繁回读验证",
        "单文件任务平均耗时：130.8 秒",
        "任务3（Python复杂）需回滚重测",
    ],
    "Qwen3.5-30B-A3（强模型）",
    [
        "显存占用：~60 GB（需 A100 级显卡）",
        "上下文窗口：32K",
        "代码生成：结构准确、逻辑清晰",
        "工具调用：准确率 ~85%+",
        "FileWrite：可一次性重写整个文件",
        "单文件任务平均耗时：77.1 秒",
        "全部单文件任务一次通过，无需回滚",
    ],
)

# Slide 4: Task 1 Python Analysis
add_two_column_slide(
    prs,
    "任务 1：Python 代码分析 — dev_cluster.py",
    "Qwen3.5-9B：149.7s ✅",
    [
        "准确识别 7 个数据类 + 4 个服务管理类",
        "完整列出 Minio / Prometheus / Grafana",
        "正确解释端口分配规则",
        "输出：详细的 Markdown 表格和架构说明",
    ],
    "Qwen3.5-30B-A3：129.0s ✅",
    [
        "分析质量与 9B 基本一致",
        "输出结构更紧凑，重点更突出",
        "耗时缩短 14%",
    ],
)

# Slide 5: Task 2 Python Simple Edit
add_two_column_slide(
    prs,
    "任务 2：Python 简单修改 — 添加 --version 参数",
    "Qwen3.5-9B：76.1s ✅",
    [
        "FileEdit 定位准确",
        "AST 验证通过",
        "需两次 FileRead 验证位置",
    ],
    "Qwen3.5-30B-A3：43.0s ✅",
    [
        "一次 FileEdit 成功",
        "无需额外验证读取",
        "耗时缩短 43%",
    ],
)

# Slide 6: Task 3 Python Complex Edit
add_two_column_slide(
    prs,
    "任务 3：Python 复杂修改 — stream_until_eof 添加时间戳",
    "Qwen3.5-9B：251.5s → 回滚重测 → 251.6s ✅",
    [
        "首次：Loop Guard 触发",
        "仅添加 import datetime",
        "未修改 stream_until_eof 函数体",
        "回滚后重测才成功",
        "总耗时超过 500 秒",
    ],
    "Qwen3.5-30B-A3：118.0s ✅ 一次通过！",
    [
        "首次即成功完成全部修改",
        "import datetime + 函数体修改",
        "stdout 和日志文件均添加时间戳",
        "无需回滚，无需重测",
        "耗时仅为 9B 的 47%",
    ],
)

# Slide 7: Task 4 C++ Analysis
add_two_column_slide(
    prs,
    "任务 4：C++ 代码分析 — cache_probe.cc",
    "Qwen3.5-9B：105.3s ✅",
    [
        "准确区分 _metrics 和 _public_metrics",
        "完整列出 15+ 个 metrics",
        "正确识别 shard 0 独占指标",
    ],
    "Qwen3.5-30B-A3：79.4s ✅",
    [
        "分析深度与 9B 一致",
        "额外读取了 cache_probe.h 头文件",
        "耗时缩短 25%",
    ],
)

# Slide 8: Task 5 C++ Simple Edit
add_two_column_slide(
    prs,
    "任务 5：C++ 简单修改 — io_result.cc 提取 to_string()",
    "Qwen3.5-9B：125.3s ✅",
    [
        "FileEdit 尝试修改函数",
        "因格式不匹配导致多次尝试",
        "最终成功提取 to_string()",
        "字符串值与原始输出有差异",
    ],
    "Qwen3.5-30B-A3：40.0s ✅",
    [
        "直接用 FileWrite 重写整个文件",
        "干净利落，一步到位",
        "耗时仅为 9B 的 32%",
    ],
)

# Slide 9: Task 6 C++ Complex Edit
add_two_column_slide(
    prs,
    "任务 6：C++ 复杂修改 — cache_probe.cc 添加 cache_hit_ratio",
    "Qwen3.5-9B：76.6s ✅",
    [
        "FileEdit 一次成功插入 gauge",
        "Lambda 计算逻辑正确",
        "除零保护完善",
    ],
    "Qwen3.5-30B-A3：53.1s ✅",
    [
        "同样一次成功",
        "代码风格与周围完全一致",
        "耗时缩短 31%",
    ],
)

# Slide 10: Task 7 Multi-file Enhancement
add_two_column_slide(
    prs,
    "任务 7：多文件功能增强 — 配置验证器 + 摘要打印",
    "Qwen3.5-9B：超时 ❌",
    [
        "Context size exceeded（上下文超限）",
        "简化 prompt 后第二次尝试",
        "超时 300s，仅完成 1/2",
        "dev_cluster.py 修改完成",
        "cluster_utils.py 未创建",
    ],
    "Qwen3.5-30B-A3：超时 ❌（部分完成）",
    [
        "未触发上下文超限",
        "超时 300s，完成度约 60%",
        "dev_cluster.py：参数 + import 完成",
        "cluster_utils.py：已创建但路径错误",
        "缺少 node_metas 后的调用逻辑",
    ],
)

# Slide 11: Summary
add_content_slide(
    prs,
    "总结与结论",
    [
        "Qwen3.5-30B-A3 在单文件任务上全面优于 9B：",
        "  • 平均耗时缩短 41%（130.8s → 77.1s）",
        "  • 任务3（Python复杂）从需回滚重测 → 一次通过",
        "  • 任务5（C++简单）从多次 FileEdit → 一次性 FileWrite",
        "  • 无一次触发 Loop Guard",
        "",
        "多文件任务仍是共同瓶颈：",
        "  • 上下文窗口（32K）对大文件+多文件场景不足",
        "  • Windows 路径解析消耗大量工具调用轮次",
        "  • 建议：拆分任务或使用 128K+ 上下文模型",
        "",
        "PilotCode 补偿机制的有效性：",
        "  • Auto-allow 和 Loop Guard 对 9B 至关重要",
        "  • 对 30B-A3 则可减少干预，提升效率",
    ],
)

output_path = r"D:\Source\2026\P2\PilotCode\tests\e2e\tasks\PilotCode_E2E_Test_Report_v2.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
